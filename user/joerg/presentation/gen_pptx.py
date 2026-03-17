#!/usr/bin/env python3
"""Generate an editable PPTX presentation for titan-oellm (Monokai theme)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Monokai Color Palette ─────────────────────────────────────
BG        = RGBColor(0x27, 0x28, 0x22)  # Monokai background
BG_LIGHT  = RGBColor(0x3E, 0x3D, 0x32)  # Monokai selection/highlight
TEXT      = RGBColor(0xF8, 0xF8, 0xF2)  # Monokai foreground
HEADING   = RGBColor(0x66, 0xD9, 0xEF)  # Monokai cyan
SUBHEAD   = RGBColor(0xAE, 0x81, 0xFF)  # Monokai purple
ACCENT    = RGBColor(0xFD, 0x97, 0x1F)  # Monokai orange
GREEN     = RGBColor(0xA6, 0xE2, 0x2E)  # Monokai green
PINK      = RGBColor(0xF9, 0x26, 0x72)  # Monokai pink/red
YELLOW    = RGBColor(0xE6, 0xDB, 0x74)  # Monokai yellow
CODE_BG   = RGBColor(0x1E, 0x1F, 0x1C)  # Darker than monokai bg
CODE_FG   = RGBColor(0xF8, 0xF8, 0xF2)  # Same as text (readable)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DIM       = RGBColor(0x75, 0x71, 0x5E)  # Monokai comment color
TABLE_HDR = RGBColor(0x49, 0x48, 0x3E)  # Monokai dark highlight
TABLE_ROW = RGBColor(0x3E, 0x3D, 0x32)  # Monokai selection

# ── Layout constants ──────────────────────────────────────────
MARGIN = Inches(0.5)
SLIDE_W = Inches(10)
CONTENT_W = Inches(9.0)
TITLE_TOP = Inches(0.3)
BODY_TOP = Inches(1.15)

MONO_FONT = "Courier New"  # Most reliable monospace across platforms


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, left=MARGIN, top=TITLE_TOP, width=CONTENT_W,
                  height=Inches(0.7), size=Pt(32), color=HEADING, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def add_text_box(slide, text="", left=MARGIN, top=BODY_TOP, width=CONTENT_W,
                 height=Inches(5.5), size=Pt(16), color=TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    return tf


def add_bullet(tf, text, level=0, size=Pt(16), color=TEXT, bold=False,
               space_before=Pt(4), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_para(tf, text, size=Pt(16), color=TEXT, bold=False, space_before=Pt(2),
             space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_code_box(slide, lines, left=MARGIN, top=Inches(2.5), width=CONTENT_W,
                 height=None, font_size=Pt(11)):
    """Add a code block as a rounded rectangle with monospace text."""
    if height is None:
        # ~18pt per line + padding
        line_h = font_size.pt * 1.35
        height = Inches((line_h * len(lines) + 16) / 72)
        height = min(height, Inches(5.5))

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.02

    tf = shape.text_frame
    tf.word_wrap = False  # No wrapping for code!
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    tf.auto_size = None

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = CODE_FG
        p.font.name = MONO_FONT
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * 1.3)
    return shape


def add_table_slide(slide, title, headers, rows, top=Inches(1.3), left=MARGIN):
    """Add a table. If title is non-empty, also adds a title box."""
    if title:
        add_title_box(slide, title)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = CONTENT_W
    row_h = Inches(0.35)
    height = row_h * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = HEADING
            p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW

    return table_shape


def first_or_add(tf, text, size=Pt(16), color=TEXT, bold=False):
    """Set text on first empty paragraph or add new one."""
    p0 = tf.paragraphs[0]
    if p0.text == "":
        p0.text = text
        p0.font.size = size
        p0.font.color.rgb = color
        p0.font.bold = bold
        return p0
    return add_bullet(tf, text, size=size, color=color, bold=bold)


# ══════════════════════════════════════════════════════════════
#  Build presentation
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ── SLIDE 1: Title ───────────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Titan-OELLM", top=Inches(2.0), size=Pt(54), color=HEADING)
tf = add_text_box(s, "", top=Inches(3.3))
p = tf.paragraphs[0]
p.text = "A TorchTitan Wrapper for Scalable LLM Training"
p.font.size = Pt(24)
p.font.color.rgb = GREEN
add_para(tf, "on HPC Systems", size=Pt(24), color=GREEN)


# ── SLIDE 2: Agenda ──────────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Agenda")
tf = add_text_box(s, "", top=Inches(1.3))
items = [
    ("1.", "Motivation", "Why a wrapper around TorchTitan?"),
    ("2.", "Architecture Overview", "Repository structure & design"),
    ("3.", "Configuration System", "TOML configs, cluster paths, environment"),
    ("4.", "Starting an Experiment", "Local & SLURM workflows"),
    ("5.", "Key Components", "Data loaders, validators, schedulers, logging"),
    ("6.", "Adding New Models", "Step-by-step guide"),
    ("7.", "Adding New Features", "Datasets, tokenizers, schedulers"),
    ("8.", "Live Demo / Q&A", ""),
]
for i, (num, title, desc) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = f"{num} "
    r.font.size = Pt(18)
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r2 = p.add_run()
    r2.text = title
    r2.font.size = Pt(18)
    r2.font.color.rgb = WHITE
    r2.font.bold = True
    if desc:
        r3 = p.add_run()
        r3.text = f"  --  {desc}"
        r3.font.size = Pt(15)
        r3.font.color.rgb = DIM
    p.space_before = Pt(8)


# ── SLIDE 3: What is TorchTitan? ────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "What is TorchTitan?")
tf = add_text_box(s, "", top=Inches(1.2))
p = tf.paragraphs[0]
r = p.add_run()
r.text = "TorchTitan (v0.2.1)"
r.font.size = Pt(18)
r.font.color.rgb = ACCENT
r.font.bold = True
r2 = p.add_run()
r2.text = " is PyTorch's native platform for large-scale LLM training."
r2.font.size = Pt(18)
r2.font.color.rgb = TEXT

add_para(tf, "", size=Pt(6))
add_para(tf, "It provides:", size=Pt(17), color=TEXT, bold=True)
for item in [
    "Multi-dimensional parallelism (FSDP2, Tensor Parallel, Pipeline Parallel, Context Parallel)",
    "Distributed checkpointing (DCP)",
    "Configuration management via TOML + tyro",
    "torch.compile integration",
    "Built-in models (Llama3, etc.)",
]:
    add_bullet(tf, item, level=0, size=Pt(16))
add_para(tf, "", size=Pt(10))
p = tf.add_paragraph()
p.text = "TorchTitan handles the distributed training engine.\nWe need everything else around it for our HPC environments."
p.font.size = Pt(16)
p.font.color.rgb = SUBHEAD
p.font.italic = True


# ── SLIDE 4: Why a Wrapper? ─────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Why a Wrapper?")
add_table_slide(s, "",
    ["Challenge", "Our Solution"],
    [
        ["Multiple HPC clusters", "cluster_config.py -- auto-detection & path resolution"],
        ["Per-user dataset/tokenizer paths", "user/$USER/cluster_paths.toml"],
        ["Custom model architectures", "Model registration system (TrainSpec)"],
        ["Validation during training", "Multi-metric validator (Perplexity, WikiText, LAMBADA)"],
        ["Training diagnostics", "Parameter & gradient logging to TensorBoard"],
        ["Flexible LR schedules", "Universal 3-phase scheduler"],
        ["Efficient data loading", "MMap, Chunked, Deterministic Packed dataloaders"],
        ["Containerized execution", "Apptainer/Singularity integration"],
    ],
    top=Inches(1.1),
)


# ── SLIDE 5: Architecture Overview ──────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Architecture Overview")
add_code_box(s, [
    "titan-oellm/",
    "  torchtitan/            <-- Git submodule (upstream, untouched)",
    "  titan_oellm/           <-- Our extensions",
    "    cluster_config.py         Cluster auto-detection & paths",
    "    configs/                  Extended JobConfig",
    "    components/               Validator, LR scheduler, param logger",
    "    datasets/                 Dataloaders, tokenizer, collator",
    "    models/                   Custom model implementations",
    "  titan_train.py         <-- Entry point (config monkey-patch)",
    "  submit_job.sh          <-- Unified job submission",
    "  slurm/                 <-- Per-cluster SLURM scripts",
    "  user/                  <-- Per-user configurations",
], top=Inches(1.2), font_size=Pt(13))
tf = add_text_box(s, "", top=Inches(5.2))
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Key principle: "
r.font.size = Pt(16)
r.font.color.rgb = TEXT
r2 = p.add_run()
r2.text = "TorchTitan is used as-is (submodule). All customization lives in titan_oellm/."
r2.font.size = Pt(16)
r2.font.color.rgb = GREEN
r2.font.bold = True


# ── SLIDE 6: Design Philosophy ───────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Design Philosophy")
# Simplified diagram - no ASCII boxes, use bullet structure instead
tf = add_text_box(s, "", top=Inches(1.2))
p = tf.paragraphs[0]
p.text = "Layered architecture:"
p.font.size = Pt(18)
p.font.color.rgb = TEXT
p.font.bold = True

add_para(tf, "", size=Pt(6))

# Layer 1
p = tf.add_paragraph()
r = p.add_run()
r.text = "Layer 1: TorchTitan Core  "
r.font.size = Pt(17)
r.font.color.rgb = HEADING
r.font.bold = True
r2 = p.add_run()
r2.text = "(used as-is, git submodule)"
r2.font.size = Pt(15)
r2.font.color.rgb = DIM
p.space_before = Pt(8)
for item in ["Parallelism (FSDP2, TP, PP, CP)", "Distributed checkpointing", "Training loop & compilation"]:
    add_bullet(tf, item, level=1, size=Pt(14), color=TEXT)

add_para(tf, "", size=Pt(6))

# Layer 2
p = tf.add_paragraph()
r = p.add_run()
r.text = "Layer 2: titan_oellm  "
r.font.size = Pt(17)
r.font.color.rgb = GREEN
r.font.bold = True
r2 = p.add_run()
r2.text = "(our extensions)"
r2.font.size = Pt(15)
r2.font.color.rgb = DIM
p.space_before = Pt(8)

components = [
    ("cluster_config", "path management"),
    ("models", "Qwen3-Custom, MoE"),
    ("datasets", "MMap, Chunked, DeterministicPacked"),
    ("components", "validator, LR scheduler, parameter logger"),
]
for name, desc in components:
    p = tf.add_paragraph()
    p.level = 1
    r = p.add_run()
    r.text = f"{name}: "
    r.font.size = Pt(14)
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT
    p.space_before = Pt(2)

add_para(tf, "", size=Pt(8))
p = tf.add_paragraph()
p.text = "All components are pluggable via the TrainSpec registration pattern."
p.font.size = Pt(16)
p.font.color.rgb = GREEN
p.font.bold = True


# ── SLIDE 7: Configuration System ────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Configuration System")
tf = add_text_box(s, "", top=Inches(1.2))
p = tf.paragraphs[0]
p.text = "Three layers of configuration:"
p.font.size = Pt(18)
p.font.color.rgb = TEXT
p.font.bold = True

for num, title, path, desc in [
    ("1.", "Model training config (TOML)", "models/qwen3_custom/train_configs/qwen3_custom.toml", "model architecture, optimizer, parallelism"),
    ("2.", "User cluster paths (TOML)", "user/$USER/cluster_paths.toml", "dataset/tokenizer paths per cluster"),
    ("3.", "Environment variables", "TITAN_USER, CLUSTER, DATASET, TOKENIZER, CONFIG, NPROC", "runtime overrides"),
]:
    add_para(tf, "", size=Pt(4))
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = f"{num} "
    r.font.size = Pt(17)
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r2 = p.add_run()
    r2.text = title
    r2.font.size = Pt(17)
    r2.font.color.rgb = WHITE
    r2.font.bold = True
    p.space_before = Pt(6)
    p2 = tf.add_paragraph()
    p2.text = f"     {desc}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = DIM
    p2.space_before = Pt(1)
    p3 = tf.add_paragraph()
    p3.text = f"     {path}"
    p3.font.size = Pt(12)
    p3.font.color.rgb = SUBHEAD
    p3.font.name = MONO_FONT
    p3.space_before = Pt(1)

add_para(tf, "", size=Pt(8))
p = tf.add_paragraph()
p.text = "Configs are cluster-independent. Paths are injected at runtime by cluster_config.py."
p.font.size = Pt(16)
p.font.color.rgb = SUBHEAD
p.font.italic = True


# ── SLIDE 8: Training Config (TOML) ─────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Training Config (TOML)")
add_code_box(s, [
    '[model]',
    'name = "qwen3_custom"',
    'flavor = "1.7B"              # Model size variant',
    '',
    '[training]',
    'steps = 100000',
    'seq_len = 2048',
    'local_batch_size = 4',
    'mixed_precision_param = "bfloat16"',
    '',
    '[optimizer]',
    'name = "AdamW"',
    'lr = 3e-4',
    '',
    '[lr_scheduler]',
    'scheduler_type = "universal"',
    'warm_steps = 1000',
    'warm_type = "linear"',
    'main_decay_type = "cosine"',
    'cooldown_steps = 500',
    '',
    '[parallelism]',
    'dp_shard = -1                # Auto',
    'tp = 1',
], top=BODY_TOP, font_size=Pt(12))


# ── SLIDE 9: User Cluster Paths ─────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "User Cluster Paths")
add_code_box(s, [
    '# user/$USER/cluster_paths.toml',
    '',
    '["cluster.juwels"]',
    'output_dir = "/p/scratch/project/user/titan_output"',
    'cache_base = "/p/scratch/project/user/cache"',
    '',
    '["tokenizer.neox.juwels"]',
    'path = "/p/scratch/project/tokenizers/neox"',
    '',
    '["dataset.slimpajama_627b.neox.juwels"]',
    'train_prefix = "/p/scratch/project/data/slimpajama/train"',
    'validation_prefix = "/p/scratch/project/data/slimpajama/val"',
    'dataloader = "ChunkedMMapDataset"',
    'min_doc_len = 128',
], top=BODY_TOP, font_size=Pt(13))
tf = add_text_box(s, "", top=Inches(5.5))
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Lookup pattern: "
r.font.size = Pt(16)
r.font.color.rgb = TEXT
r2 = p.add_run()
r2.text = "dataset.{name}.{tokenizer}.{cluster}"
r2.font.size = Pt(16)
r2.font.color.rgb = GREEN
r2.font.name = MONO_FONT
r2.font.bold = True


# ── SLIDE 10: Path Resolution Flow ──────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Path Resolution Flow")
add_code_box(s, [
    "Environment Variables         cluster_paths.toml",
    "+------------------+         +------------------+",
    "| DATASET=slimpaj  |         | dataset paths    |",
    "| TOKENIZER=neox   |---+     | tokenizer paths  |",
    "| CLUSTER=juwels   |   |     | cluster settings |",
    "+------------------+   |     +--------+---------+",
    "                       |              |",
    "                  +----v--------------v----+",
    "                  |   cluster_config.py    |",
    "                  |   get_cli_args()       |",
    "                  +------------+-----------+",
    "                               |",
    "                  +------------v-----------+",
    "                  | --job.config_file=...  |",
    "                  | --data.data_prefix=... |",
    "                  | --data.chunks_dir=...  |",
    "                  | --validation.data_pref |",
    "                  | --benchmarks.wikitext2 |",
    "                  +------------------------+",
    "                     CLI args for torchrun",
], top=BODY_TOP, font_size=Pt(11))


# ── SLIDE 11: Starting an Experiment ─────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Starting an Experiment")
tf = add_text_box(s, "", top=BODY_TOP)
p = tf.paragraphs[0]
p.text = "Option 1: Local (development/testing)"
p.font.size = Pt(18)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    'TITAN_USER=joerg DATASET=test_dataset \\',
    'CONFIG=user/joerg/configs/debug.toml NPROC=1 \\',
    '  bash submit_job.sh --local \\',
    '    --model.flavor=debugmodel --training.seq_len=512',
], top=Inches(1.7), font_size=Pt(13))
tf2 = add_text_box(s, "", top=Inches(3.5))
p = tf2.paragraphs[0]
p.text = "Option 2: SLURM (production)"
p.font.size = Pt(18)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    'TITAN_USER=joerg DATASET=slimpajama_627b \\',
    '  bash submit_job.sh --nodes=4 -- \\',
    '    --model.flavor=1.7B --training.steps=100000',
], top=Inches(4.1), font_size=Pt(13))


# ── SLIDE 12: Experiment Launch Flow ─────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Experiment Launch Flow")
add_code_box(s, [
    "submit_job.sh",
    "  |",
    "  +-- [--local mode]",
    "  |     +-- Load cluster_config.py",
    "  |     +-- Resolve paths  (get_cli_args)",
    "  |     +-- Set env vars   (get_env_exports)",
    "  |     +-- apptainer exec ... torchrun -m torchtitan.train",
    "  |",
    "  +-- [SLURM mode]",
    "        +-- Create .venv_submit (no torch on login node)",
    "        +-- Detect cluster from hostname",
    "        +-- Find container  (titan_${CLUSTER}_0.2.1.sif)",
    "        +-- sbatch slurm/$CLUSTER.sh $TRAINING_ARGS",
    "              +-- Load modules (CUDA, NCCL, ...)",
    "              +-- Load cluster_config (inside container)",
    "              +-- Set NCCL env vars",
    "              +-- torchrun --nnodes=$N -m torchtitan.train",
], top=BODY_TOP, font_size=Pt(13))


# ── SLIDE 13: Data Loaders Overview ──────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Key Components: Data Loaders")
add_table_slide(s, "",
    ["Type", "Use Case", "Key Feature"],
    [
        ["MMapDataset", "Single large binary file", "Memory-mapped random access"],
        ["ChunkedMMapDataset", "Pre-chunked data", "Deterministic, validation split"],
        ["DeterministicPackedDataset", "Document packing", "Fixed-length seqs, reproducible"],
    ],
    top=Inches(1.2),
)
add_code_box(s, [
    "Raw Text --> Tokenize --> .bin/.idx files --> MMap/Chunked Loader",
    "                                                   |",
    "                                             collate_function()",
    "                                                   |",
    "                                             (input_ids, labels)",
], top=Inches(3.2), font_size=Pt(12))
tf = add_text_box(s, "", top=Inches(5.0))
p = tf.paragraphs[0]
p.text = "Configured via [data] section and cluster_paths.toml."
p.font.size = Pt(15)
p.font.color.rgb = DIM


# ── SLIDE 14: Why Chunks? ───────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Why Chunks? The Problem with Single MMap Files")
tf = add_text_box(s, "", top=BODY_TOP)
p = tf.paragraphs[0]
p.text = "Scaling issues with one big .bin file:"
p.font.size = Pt(17)
p.font.color.rgb = PINK
p.font.bold = True

for title, desc in [
    ("Shuffling", "Sequential file + random access across 2TB = I/O thrashing on GPFS/Lustre"),
    ("Multi-node I/O", "All ranks read from one file = contention"),
    ("Validation split", "Hard to hold out data without a second copy"),
    ("Reproducibility", "Different dp_world_size changes which data each rank sees"),
]:
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = f"{title}: "
    r.font.size = Pt(15)
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(15)
    r2.font.color.rgb = TEXT
    p.space_before = Pt(6)

add_para(tf, "", size=Pt(8))
p = tf.add_paragraph()
p.text = "What chunks solve:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
p.space_before = Pt(8)

for item in [
    "Pre-shuffled chunks (e.g. 256 x ~8GB) -- sequential reads are already random",
    "Round-robin assignment -- each DP rank gets its own chunks, no contention",
    "Built-in validation split -- reserve first N docs/chunk for validation",
    "Deterministic across node counts -- same seed + same chunks = same order",
]:
    add_bullet(tf, item, size=Pt(14))


# ── SLIDE 15: preprocess_mmap_chunks ─────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "preprocess_mmap_chunks: Creating Chunks")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "titan_oellm/datasets/utils/preprocess_mmap_chunks.py"
p.font.size = Pt(12)
p.font.color.rgb = DIM
p.font.name = MONO_FONT

add_code_box(s, [
    "Input:  train.bin (2TB) + train.idx",
    "              |",
    "  preprocess_mmap_chunks.py",
    "  (parallel workers + async flush)",
    "              |",
    "Output: chunks_dir/",
    "          chunk_0000.bin + chunk_0000.idx",
    "          chunk_0001.bin + chunk_0001.idx",
    "          ...           ",
    "          chunk_0255.bin + chunk_0255.idx",
], top=Inches(1.5), font_size=Pt(13))

tf2 = add_text_box(s, "", top=Inches(4.3))
p = tf2.paragraphs[0]
p.text = "Key features:"
p.font.size = Pt(16)
p.font.color.rgb = GREEN
p.font.bold = True
for item in [
    "Multi-process: ProcessPoolExecutor for parallel reading",
    "Async disk writes: separate writer processes so readers never block",
    "Random assignment: each document randomly assigned to a target chunk",
    "Per-chunk shuffle: each chunk shuffled independently after assignment",
    "Validation + cleanup: automatic integrity check, temp files deleted after",
]:
    add_bullet(tf2, item, size=Pt(14))


# ── SLIDE 16: ChunkedMMapDataset ─────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "ChunkedMMapDataset: How It Works")
add_code_box(s, [
    "      All chunks (sorted, deterministic)",
    "      [chunk_00, chunk_01, chunk_02, ..., chunk_N]",
    "                |",
    "      Seed-based permutation per epoch",
    "                |",
    "      Round-robin assignment to DP ranks",
    "                |",
    "  Rank 0            Rank 1            Rank 2",
    "  [chunk_02,        [chunk_07,        [chunk_00,",
    "   chunk_05,         chunk_01,         chunk_09,",
    "   chunk_11, ...]    chunk_04, ...]    chunk_06, ...]",
    "      |                  |                  |",
    "  Sequential read   Sequential read   Sequential read",
    "  (fast, no seeks)  (fast, no seeks)  (fast, no seeks)",
], top=BODY_TOP, font_size=Pt(12))

tf = add_text_box(s, "", top=Inches(5.0))
for item in [
    "Chunks are pre-shuffled at creation time, so sequential = random",
    "New epoch = re-shuffle chunk order (seed + epoch_counter)",
    "Validation split: use_only_first_n_per_chunk / exclude_first_n_per_chunk",
]:
    first_or_add(tf, item, size=Pt(14))


# ── SLIDE 17: DeterministicPackedDataset Deep Dive ───────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "DeterministicPackedDataset: Deep Dive")
tf = add_text_box(s, "", top=BODY_TOP)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "The problem: "
r.font.size = Pt(16)
r.font.color.rgb = PINK
r.font.bold = True
r2 = p.add_run()
r2.text = "Documents have variable length. Padding wastes compute. Per-rank streaming is non-deterministic across node counts."
r2.font.size = Pt(16)
r2.font.color.rgb = TEXT

add_para(tf, "", size=Pt(4))
p = tf.add_paragraph()
p.text = "Greedy packing approach:"
p.font.size = Pt(16)
p.font.color.rgb = GREEN
p.font.bold = True

add_code_box(s, [
    "Doc1 (500) | Doc2 (1200) | Doc3 (800) | Doc4 (300) | ...",
    "                    |",
    "      Treat as one continuous token stream",
    "                    |",
    "  [-- seq_len+1 --][-- seq_len+1 --][-- seq_len+1 --]",
    "     Sequence 0       Sequence 1       Sequence 2",
], top=Inches(2.6), font_size=Pt(13))

tf2 = add_text_box(s, "", top=Inches(4.3))
for item in [
    "Documents concatenated (+ optional EOS) into a virtual token stream",
    "Fixed-length sequences cut at regular positions",
    "No stored index -- binary search on cumulative token counts (32KB metadata)",
    "Checkpoint = 1 integer (global_sequence_id) -- instant resume",
]:
    first_or_add(tf2, item, size=Pt(14))


# ── SLIDE 18: Batch Diversity ────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "DeterministicPackedDataset: Batch Diversity")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "Strided assignment: every batch samples across the whole dataset"
p.font.size = Pt(16)
p.font.color.rgb = TEXT

add_code_box(s, [
    "Epoch token stream (all chunks, permuted):",
    "[=========================================================]",
    " ^          ^          ^          ^          ^",
    " lane 0    lane 1    lane 2    lane 3    lane 4",
    " (step 0)  (step 0)  (step 0)  (step 0)  (step 0)",
    "",
    "Step 0 batch = {lane_0[0], lane_1[0], lane_2[0], ...}",
    "Step 1 batch = {lane_0[1], lane_1[1], lane_2[1], ...}",
], top=Inches(1.7), font_size=Pt(13))

tf2 = add_text_box(s, "", top=Inches(4.0))
for item in [
    "Each lane advances sequentially through its region (fast I/O)",
    "Each batch spans the entire dataset (diversity)",
    "Batch composition depends only on global_batch_size, NOT on dp_world_size",
    "Checkpoint = 1 integer (global_sequence_id) -- instant resume",
]:
    first_or_add(tf2, item, size=Pt(15))

add_para(tf2, "", size=Pt(6))
p = tf2.add_paragraph()
p.text = "Scales to: 1-10T tokens, 100-4000 GPUs, global_batch_size up to 16M."
p.font.size = Pt(16)
p.font.color.rgb = GREEN
p.font.bold = True


# ── SLIDE 19: Validator ──────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Key Components: Validator")
tf = add_text_box(s, "", top=BODY_TOP)
p = tf.paragraphs[0]
p.text = "Multi-metric validation during training:"
p.font.size = Pt(17)
p.font.color.rgb = TEXT
p.font.bold = True

for metric, desc in [
    ("Perplexity", "cross-entropy loss on validation set"),
    ("WikiText-2 / WikiText-103", "standard LM benchmarks"),
    ("LAMBADA", "last-word prediction accuracy"),
    ("Spearman Correlation", "rank correlation metric"),
]:
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = f"{metric}: "
    r.font.size = Pt(16)
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(16)
    r2.font.color.rgb = TEXT
    p.space_before = Pt(6)

add_code_box(s, [
    '[validation]',
    'enable = true',
    'freq = 1000          # Validate every N steps',
    'eval_mode = "concatenated"  # or "document"',
    '',
    '[benchmarks]',
    'wikitext2_path = "..."    # Injected by cluster_config',
    'lambada_path = "..."',
], top=Inches(3.8), font_size=Pt(12))
tf2 = add_text_box(s, "", top=Inches(5.9))
p = tf2.paragraphs[0]
p.text = "Supports multiple validation datasets per training run."
p.font.size = Pt(15)
p.font.color.rgb = SUBHEAD
p.font.bold = True


# ── SLIDE 20: Universal LR Scheduler ────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Key Components: Universal LR Scheduler")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "Three-phase learning rate schedule:"
p.font.size = Pt(17)
p.font.color.rgb = TEXT
add_code_box(s, [
    " LR",
    "  ^",
    "  |    Phase 1      Phase 2         Phase 3",
    "  |   (Warm)        (Main)         (Cooldown)",
    "  |",
    "  |      /```````````````````\\",
    "  |     /                     \\",
    "  |    /  linear/cosine/       \\ cosine/linear",
    "  |   /   constant decay        \\ decay",
    "  |  /                           \\___________",
    "  | /                                 lr_min",
    "  +-+-------+------------------+-------+----->",
    "    0   warm_steps        total-cooldown  steps",
], top=Inches(1.6), font_size=Pt(12))
tf2 = add_text_box(s, "", top=Inches(5.0))
for item in [
    "Supports: linear, cosine, sqrt, exp decay types per phase",
    "Warm phase can go up (warmup) or down (warmdown for continued pretraining)",
    "Full checkpoint/resume support",
]:
    first_or_add(tf2, item, size=Pt(15))


# ── SLIDE 21: Parameter Logger ───────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Key Components: Parameter Logger")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "Statistical insights into training dynamics:"
p.font.size = Pt(17)
p.font.color.rgb = TEXT
add_code_box(s, [
    '[parameter_logging]',
    'enabled = true',
    'log_interval = 1000',
    'log_parameters = true       # Weight stats (max, min, norm, std)',
    'log_gradients = true        # Gradient statistics',
    'log_optimizer_states = true  # Adam momentum/variance stats',
], top=Inches(1.7), font_size=Pt(13))
tf2 = add_text_box(s, "", top=Inches(3.5))
for item in [
    "Pattern-based filtering (include/exclude specific layers)",
    "All stats logged to TensorBoard",
    "Helps diagnose: gradient explosion, dead layers, optimizer issues",
]:
    first_or_add(tf2, item, size=Pt(16))


# ── SLIDE 22: TrainSpec Pattern ──────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Model Registration: TrainSpec Pattern")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "Every model registers a TrainSpec -- wires everything together:"
p.font.size = Pt(16)
p.font.color.rgb = TEXT
add_code_box(s, [
    "TrainSpec(",
    "    model_cls       = Qwen3Model,",
    "    model_args      = qwen3_custom_configs,",
    "    parallelize_fn  = parallelize_qwen3,",
    "    pipelining_fn   = ...,",
    "    build_optimizers_fn     = ...,",
    "    build_lr_schedulers_fn  = ...,",
    "    build_dataloader_fn     = build_sci_dataloader,",
    "    build_tokenizer_fn      = build_sci_hf_tokenizer,",
    "    build_loss_fn           = ...,",
    "    build_validator_fn      = build_validator,",
    "    build_metrics_processor_fn = ...,",
    "    state_dict_adapter      = Qwen3StateDictAdapter,",
    ")",
    "",
    'register_train_spec("qwen3_custom", spec)',
], top=Inches(1.5), font_size=Pt(13))


# ── SLIDE 23: Qwen3-Custom Model ────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Current Model: Qwen3-Custom")
add_table_slide(s, "",
    ["Flavor", "dim", "layers", "heads", "kv_heads"],
    [
        ["debugmodel", "128", "2", "2", "2"],
        ["0.5B", "896", "24", "14", "2"],
        ["0.6B", "1024", "28", "16", "8"],
        ["1.7B", "1536", "28", "12", "2"],
        ["4B", "2560", "36", "20", "4"],
        ["8B", "3584", "36", "28", "4"],
        ["14B", "5120", "40", "40", "8"],
        ["32B", "5120", "64", "40", "8"],
    ],
    top=Inches(1.2),
)
tf = add_text_box(s, "", top=Inches(4.9))
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Plus: "
r.font.size = Pt(14)
r.font.color.rgb = TEXT
r2 = p.add_run()
r2.text = "125M, 125M768, 130Msci, 1.7Bsci "
r2.font.size = Pt(14)
r2.font.color.rgb = ACCENT
r2.font.bold = True
r3 = p.add_run()
r3.text = "(custom research sizes) and "
r3.font.size = Pt(14)
r3.font.color.rgb = TEXT
r4 = p.add_run()
r4.text = "MoE variants "
r4.font.size = Pt(14)
r4.font.color.rgb = ACCENT
r4.font.bold = True
r5 = p.add_run()
r5.text = "(debugmodel_moe, 600M-A60M)"
r5.font.size = Pt(14)
r5.font.color.rgb = TEXT
add_para(tf, "Features: QK-norm, configurable RoPE theta, HF weight loading, depth init, weight tying.",
         size=Pt(14), color=DIM)


# ── SLIDE 24: Adding a New Model (1) ────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Adding a New Model -- Step by Step")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "1. Create the directory structure:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    "titan_oellm/models/my_model/",
    "  __init__.py              # TrainSpec registration",
    "  model/",
    "    args.py                # Model hyperparameters",
    "    model.py               # nn.Module implementation",
    "    state_dict_adapter.py  # (optional) HF weight conversion",
    "  infra/",
    "    parallelize.py         # FSDP/TP/AC setup",
    "  train_configs/",
    "    my_model.toml          # Default training config",
], top=Inches(1.6), font_size=Pt(13))

tf2 = add_text_box(s, "", top=Inches(4.3))
p = tf2.paragraphs[0]
p.text = "2. Define model arguments:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    "@dataclass",
    "class MyModelArgs(BaseModelArgs):",
    "    dim: int = 1024",
    "    n_layers: int = 24",
    "    n_heads: int = 16",
    "    vocab_size: int = 50432",
], top=Inches(4.9), font_size=Pt(13))


# ── SLIDE 25: Adding a New Model (2) ────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Adding a New Model -- Step by Step (cont.)")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "3. Implement the model (nn.Module with forward pass)"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True

add_para(tf, "", size=Pt(4))
p = tf.add_paragraph()
p.text = "4. Implement parallelization:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    "def parallelize_my_model(model, parallel_dims, job_config):",
    "    if parallel_dims.tp_enabled:",
    "        apply_tp(model, ...)",
    "    if parallel_dims.dp_shard_enabled:",
    "        apply_fsdp(model, ...)",
    "    if job_config.activation_checkpoint.mode != 'none':",
    "        apply_ac(model, ...)",
], top=Inches(2.3), font_size=Pt(13))

tf2 = add_text_box(s, "", top=Inches(4.2))
p = tf2.paragraphs[0]
p.text = "5. Register the TrainSpec:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    'configs = {"base": MyModelArgs(...), "large": MyModelArgs(...)}',
    "spec = TrainSpec(model_cls=MyModel, model_args=configs, ...)",
    'register_train_spec("my_model", spec)',
], top=Inches(4.8), font_size=Pt(13))

tf3 = add_text_box(s, "", top=Inches(5.8))
p = tf3.paragraphs[0]
p.text = "6. Add import in titan_oellm/models/__init__.py:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    "from . import my_model   # Auto-registers on import",
], top=Inches(6.35), height=Inches(0.4), font_size=Pt(13))


# ── SLIDE 26: Adding a New Dataset ───────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Adding a New Dataset")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "1. Tokenize your data:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    "python titan_oellm/scripts/convert_hf_to_mmap.py \\",
    "    --input your_data/ --output /scratch/data/my_dataset/ \\",
    "    --tokenizer /path/to/tokenizer",
], top=Inches(1.6), font_size=Pt(13))

tf2 = add_text_box(s, "", top=Inches(2.8))
p = tf2.paragraphs[0]
p.text = "2. Register in user/$USER/cluster_paths.toml:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    '["dataset.my_dataset.neox.juwels"]',
    'train_prefix = "/scratch/data/my_dataset/train"',
    'validation_prefix = "/scratch/data/my_dataset/val"',
    'dataloader = "MMapDataset"',
    'min_doc_len = 64',
], top=Inches(3.4), font_size=Pt(13))

tf3 = add_text_box(s, "", top=Inches(5.1))
p = tf3.paragraphs[0]
p.text = "3. Use it:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    "DATASET=my_dataset bash submit_job.sh --local",
], top=Inches(5.6), height=Inches(0.4), font_size=Pt(14))


# ── SLIDE 27: Adding a Tokenizer ─────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Adding a New Tokenizer")
tf = add_text_box(s, "", top=Inches(1.05))
p = tf.paragraphs[0]
p.text = "1. Register in cluster paths:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    '["tokenizer.my_tokenizer.juwels"]',
    'path = "/scratch/tokenizers/my_tokenizer"',
], top=Inches(1.6), font_size=Pt(13))

tf2 = add_text_box(s, "", top=Inches(2.6))
p = tf2.paragraphs[0]
p.text = "2. Create dataset entries matching the new tokenizer:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    '["dataset.my_data.my_tokenizer.juwels"]',
    'train_prefix = "/scratch/data/my_data_mytok/train"',
    'validation_prefix = "/scratch/data/my_data_mytok/val"',
    'dataloader = "ChunkedMMapDataset"',
], top=Inches(3.2), font_size=Pt(13))

tf3 = add_text_box(s, "", top=Inches(4.6))
p = tf3.paragraphs[0]
p.text = "3. Use it:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
add_code_box(s, [
    "DATASET=my_data TOKENIZER=my_tokenizer bash submit_job.sh --local",
], top=Inches(5.1), height=Inches(0.4), font_size=Pt(14))


# ── SLIDE 28: Multi-Cluster Support ──────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Multi-Cluster Support")
add_table_slide(s, "",
    ["Cluster", "Location", "Detection Pattern", "Container"],
    [
        ["JUWELS", "FZ Juelich", "jwlogin*, jwc*, juwels", "titan_juwels_0.2.1.sif"],
        ["Jupiter", "FZ Juelich", "jupiter*, jrc*", "titan_jupiter_0.2.1.sif"],
        ["Capella", "PSNC", "c + digit, capella", "titan_capella_0.2.1.sif"],
        ["Leonardo", "CINECA", "leonardo", "titan_leonardo_0.2.1.sif"],
        ["Local", "Dev machine", "fallback", "local container"],
    ],
    top=Inches(1.2),
)
tf = add_text_box(s, "", top=Inches(3.8))
for item in [
    "Cluster is auto-detected from hostname. Override with CLUSTER=...",
    "Same training config runs everywhere -- only paths differ.",
]:
    first_or_add(tf, item, size=Pt(16))


# ── SLIDE 29: Execution Environment ─────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Execution Environment")
add_code_box(s, [
    "Login Node                       Compute Node",
    "+------------------+            +----------------------------+",
    "|                  |            | Apptainer Container        |",
    "| submit_job.sh    |            | +------------------------+ |",
    "| .venv_submit     |  sbatch    | | torchrun               | |",
    "| (no torch!)      | ---------> | | torchtitan.train       | |",
    "| cluster_config   |            | | titan_oellm (models,   | |",
    "|                  |            | |   data, components)    | |",
    "+------------------+            | +------------------------+ |",
    "                                | GPU 0  GPU 1  GPU 2  GPU 3|",
    "                                +----------------------------+",
], top=BODY_TOP, font_size=Pt(12))
tf = add_text_box(s, "", top=Inches(4.5))
for item in [
    "Login node: lightweight venv (no PyTorch) for config resolution only",
    "Compute node: full environment inside Apptainer container",
    "All paths bind-mounted into container",
]:
    first_or_add(tf, item, size=Pt(16))


# ── SLIDE 30: Typical Workflow ───────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Typical Workflow")
tf = add_text_box(s, "", top=Inches(1.2))
steps = [
    ("1.", "Clone repo, set up user config", "user/$USER/cluster_paths.toml"),
    ("2.", "Choose or create training config", "user/$USER/configs/my_experiment.toml"),
    ("3.", "(Optional) Tokenize data & download benchmarks", "convert_hf_to_mmap.py, download_benchmarks.py"),
    ("4.", "Submit experiment", "TITAN_USER=$USER DATASET=X bash submit_job.sh --nodes=8 -- --model.flavor=4B"),
    ("5.", "Monitor training", "tensorboard --logdir outputs/"),
]
for i, (num, step, detail) in enumerate(steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = f"{num} "
    r.font.size = Pt(17)
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r2 = p.add_run()
    r2.text = step
    r2.font.size = Pt(17)
    r2.font.color.rgb = WHITE
    r2.font.bold = True
    p.space_before = Pt(12)
    p2 = tf.add_paragraph()
    p2.text = f"     {detail}"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DIM
    p2.font.name = MONO_FONT
    p2.space_before = Pt(2)


# ── SLIDE 31: Summary ───────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Summary")
tf = add_text_box(s, "", top=BODY_TOP)
p = tf.paragraphs[0]
p.text = "What Titan-OELLM gives you:"
p.font.size = Pt(18)
p.font.color.rgb = GREEN
p.font.bold = True
for item in [
    "Seamless multi-cluster support -- same config, any HPC system",
    "Pluggable architecture -- swap models, dataloaders, schedulers",
    "Production-ready training -- validation, logging, checkpointing",
    "Easy onboarding -- user directory, environment variables, TOML configs",
    "Extensibility -- add a model in ~5 files, a dataset in ~3 lines",
]:
    add_bullet(tf, item, size=Pt(15))

add_para(tf, "", size=Pt(10))
p = tf.add_paragraph()
p.text = "What TorchTitan gives us (under the hood):"
p.font.size = Pt(18)
p.font.color.rgb = SUBHEAD
p.font.bold = True
p.space_before = Pt(8)
for item in [
    "FSDP2, Tensor Parallel, Pipeline Parallel",
    "Distributed checkpointing",
    "Training loop, optimizer, compilation",
    "All the hard distributed systems work",
]:
    add_bullet(tf, item, size=Pt(15))


# ── SLIDE 32: Questions ──────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Questions?", top=Inches(2.5), size=Pt(48))
tf = add_text_box(s, "", top=Inches(3.8))
p = tf.paragraphs[0]
p.text = "Key resources:"
p.font.size = Pt(18)
p.font.color.rgb = SUBHEAD
p.font.bold = True
for item in [
    "README.md -- Quick start & overview",
    "titan_oellm/configs/README.md -- Config system docs",
    "titan_oellm/models/qwen3_custom/README.md -- Model docs",
    "user/example/ -- Template configurations",
]:
    add_bullet(tf, item, size=Pt(16))


# ── SLIDE 33: Appendix title ────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Appendix", top=Inches(3.0), size=Pt(48))


# ── SLIDE 34: Appendix A ────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_table_slide(s, "Appendix A: Full Config Sections",
    ["TOML Section", "Purpose"],
    [
        ["[job]", "Output folder, config module"],
        ["[model]", "Architecture name, flavor, vocab size"],
        ["[training]", "Steps, batch size, seq_len, precision"],
        ["[optimizer]", "AdamW params (lr, betas, weight_decay)"],
        ["[lr_scheduler]", "Scheduler type, phase config"],
        ["[parallelism]", "DP, TP, PP, CP degrees"],
        ["[data]", "Dataloader type, data paths"],
        ["[validation]", "Enable, frequency, metrics"],
        ["[benchmarks]", "WikiText, LAMBADA paths"],
        ["[parameter_logging]", "Stats logging config"],
        ["[checkpoint]", "Enable, interval, HF loading"],
        ["[compile]", "torch.compile settings"],
        ["[activation_checkpoint]", "AC mode (full/selective)"],
    ],
)


# ── SLIDE 35: Appendix B ────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_table_slide(s, "Appendix B: Environment Variables Reference",
    ["Variable", "Default", "Description"],
    [
        ["TITAN_USER", "(required)", "Your username"],
        ["CLUSTER", "auto-detected", "Target cluster"],
        ["DATASET", "test_dataset (local)", "Dataset name"],
        ["TOKENIZER", "neox", "Tokenizer name"],
        ["CONFIG", "debug.toml (local)", "Config file"],
        ["NPROC", "1", "GPUs (local mode)"],
        ["OUTPUT_DIR", "from cluster config", "Output base path"],
    ],
)


# ── SLIDE 36: Appendix C ────────────────────────────────────

s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title_box(s, "Appendix C: HuggingFace Weight Loading")
tf = add_text_box(s, "", top=BODY_TOP)
p = tf.paragraphs[0]
p.text = "Load pretrained HF checkpoints directly:"
p.font.size = Pt(17)
p.font.color.rgb = TEXT
add_code_box(s, [
    '[checkpoint]',
    'enable_checkpoint = true',
    'load_hf_model_weights_only = true',
    'hf_model_weights_path = "/path/to/hf/model"',
], top=Inches(1.6), font_size=Pt(14))

tf2 = add_text_box(s, "", top=Inches(3.1))
p = tf2.paragraphs[0]
p.text = "The StateDictAdapter handles:"
p.font.size = Pt(17)
p.font.color.rgb = GREEN
p.font.bold = True
for item in [
    "Weight name mapping (HF naming -> TorchTitan naming)",
    "Shape conversions (e.g., fused QKV -> separate Q, K, V)",
    "Skipping incompatible layers",
]:
    add_bullet(tf2, item, size=Pt(16))
add_para(tf2, "", size=Pt(8))
p = tf2.add_paragraph()
p.text = "Supports both fine-tuning and continual pretraining."
p.font.size = Pt(16)
p.font.color.rgb = SUBHEAD
p.font.bold = True


# ── Save ─────────────────────────────────────────────────────

out_path = "/home/joerg/workspace/python/github/titan-oellm/user/joerg/presentation_titan_oellm.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
